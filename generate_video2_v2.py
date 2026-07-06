#!/usr/bin/env python3
"""
Generate UKPSC Decoded Video 2 PPTX — Note-Making Masterclass
"The Note-Making Architecture Toppers Hide"
Style: Authority Masterclass Dark-Mode
Uses python-pptx for Keynote-compatible output.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy

# --- Color Palette ---
CARBON = RGBColor(0x17, 0x17, 0x17)
CARD_BG = RGBColor(0x22, 0x22, 0x22)
CARD_BORDER = RGBColor(0x33, 0x33, 0x33)
GOLD = RGBColor(0xD4, 0xA3, 0x73)
TEAL = RGBColor(0x7B, 0x9E, 0x89)
BODY_TEXT = RGBColor(0xE5, 0xE5, 0xE5)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Fonts
TITLE_FONT = "Garamond"
BODY_FONT = "Inter"
HEADER_FONT = "Inter"


def set_slide_bg(slide, color=CARBON):
    """Set solid background color for a slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BODY_TEXT, align=PP_ALIGN.LEFT,
                font_name=BODY_FONT):
    """Add a simple single-line textbox"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT):
    """Add multiline textbox. lines = list of (text, color, bold, size, font_name) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        text = line[0]
        color = line[1] if len(line) > 1 else BODY_TEXT
        bold = line[2] if len(line) > 2 else False
        size = line[3] if len(line) > 3 else 14
        font = line[4] if len(line) > 4 else BODY_FONT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return txBox


def add_card(slide, left, top, width, height, fill=CARD_BG, border_color=None):
    """Add a rounded rectangle card"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_insider_tip(slide, left, top, width, height, tip_text):
    """Add the golden-bordered Insider Tip box"""
    card = add_card(slide, left, top, width, height, fill=CARBON, border_color=GOLD)
    txBox = slide.shapes.add_textbox(
        int(left + Inches(0.2)), int(top + Inches(0.1)),
        int(width - Inches(0.4)), int(height - Inches(0.2)))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    # Star
    r1 = p.add_run()
    r1.text = "★ "
    r1.font.size = Pt(13)
    r1.font.color.rgb = GOLD
    r1.font.name = BODY_FONT
    # "Insider Tip"
    r2 = p.add_run()
    r2.text = "Insider Tip  "
    r2.font.size = Pt(13)
    r2.font.bold = True
    r2.font.color.rgb = GOLD
    r2.font.name = BODY_FONT
    # Tip text
    r3 = p.add_run()
    r3.text = tip_text
    r3.font.size = Pt(13)
    r3.font.color.rgb = BODY_TEXT
    r3.font.name = BODY_FONT



def add_gold_line(slide, left, top, width):
    """Add a gold horizontal separator line"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


# ============================================================
# SLIDE BUILDERS
# ============================================================

def build_slide_01(prs):
    """Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    # Header
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(6.5), Inches(0.4),
                "M A S T E R C L A S S", font_size=12, bold=True, color=MUTED, font_name=HEADER_FONT)
    # Title line 1
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(6.5), Inches(0.8),
                "The Note-Making", font_size=44, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    # Title line 2
    add_textbox(slide, Inches(0.8), Inches(2.6), Inches(6.5), Inches(0.8),
                "Architecture Toppers Hide", font_size=44, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    # Gold line
    add_gold_line(slide, Inches(0.8), Inches(3.8), Inches(2.5))
    # Subtitle
    add_textbox(slide, Inches(0.8), Inches(4.1), Inches(6.5), Inches(0.5),
                "The 5-Layer System Behind 800+ Mains Scores", font_size=16, color=BODY_TEXT)
    # Three pill boxes
    pills = [("5 Layers", 0.8), ("4 Secrets", 3.0), ("1 Page Rule", 5.2)]
    for label, px in pills:
        add_card(slide, Inches(px), Inches(5.0), Inches(1.8), Inches(0.45), fill=CARBON, border_color=GOLD)
        add_textbox(slide, Inches(px), Inches(5.02), Inches(1.8), Inches(0.4),
                    label, font_size=12, color=GOLD, align=PP_ALIGN.CENTER)


def build_slide_02(prs):
    """The Diagnosis — Summary ≠ Notes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "T H E  D I A G N O S I S", font_size=12, bold=True, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.8),
                "Summary ≠ Notes", font_size=40, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.75), Inches(2.5))
    # Quote
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(6.8), Inches(0.6),
                "\"If your notes are a smaller version of the book, you haven't made notes.\"",
                font_size=15, color=GOLD, font_name=TITLE_FONT)
    # Left card — Wrong
    add_card(slide, Inches(0.8), Inches(2.9), Inches(3.2), Inches(3.0), border_color=CARD_BORDER)
    add_multiline(slide, Inches(1.0), Inches(3.0), Inches(2.9), Inches(2.8), [
        ("THE SUMMARY TRAP", MUTED, True, 11),
        ("", BODY_TEXT, False, 6),
        ("1. Read chapter 3 times", BODY_TEXT, False, 12),
        ("2. Condense to 10 pages", BODY_TEXT, False, 12),
        ("3. Re-read like a book", BODY_TEXT, False, 12),
        ("4. Same recall lag", BODY_TEXT, False, 12),
        ("", BODY_TEXT, False, 6),
        ("Result: Page count changed.", MUTED, False, 11),
        ("Processing speed didn't.", MUTED, False, 11),
    ])
    # Right card — Correct
    add_card(slide, Inches(4.3), Inches(2.9), Inches(3.2), Inches(3.0), border_color=TEAL)
    add_multiline(slide, Inches(4.5), Inches(3.0), Inches(2.9), Inches(2.8), [
        ("THE SYSTEM", TEAL, True, 11),
        ("", BODY_TEXT, False, 6),
        ("1. Demand-deconstruct PYQs", BODY_TEXT, False, 12),
        ("2. Remove overlap", BODY_TEXT, False, 12),
        ("3. Structure by dimension", BODY_TEXT, False, 12),
        ("4. Instant retrieval", BODY_TEXT, False, 12),
        ("", BODY_TEXT, False, 6),
        ("Result: Recall is instant.", TEAL, False, 11),
        ("Structure = muscle memory.", TEAL, False, 11),
    ])
    add_insider_tip(slide, Inches(0.8), Inches(6.2), Inches(6.8), Inches(0.55),
                    "Compression without a framework doesn't create recall pathways.")



def build_slide_03(prs):
    """Layer 1 — Content Sourcing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "L A Y E R  1", font_size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "Content Sourcing", font_size=38, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.7), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(6.5), Inches(0.4),
                "Three sources. One system. Everything else is noise.", font_size=14, color=MUTED)
    # Three source cards
    sources = [
        ("Syllabus", "The skeleton. Every note\nmust trace back to a\nsyllabus line."),
        ("PYQs", "The demand map. What\nexams actually ask —\nnot what you assume."),
        ("Current Affairs", "The living layer. Folded\ninto structure, never\nstored separately."),
    ]
    for i, (title, desc) in enumerate(sources):
        x = Inches(0.8 + i * 2.4)
        add_card(slide, x, Inches(2.6), Inches(2.1), Inches(2.8), border_color=CARD_BORDER)
        add_textbox(slide, x, Inches(2.75), Inches(2.1), Inches(0.4),
                    title, font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_multiline(slide, int(x + Inches(0.15)), Inches(3.3), Inches(1.8), Inches(1.8),
                      [(line, BODY_TEXT, False, 11) for line in desc.split("\n")],
                      align=PP_ALIGN.CENTER)
    add_insider_tip(slide, Inches(0.8), Inches(5.8), Inches(6.8), Inches(0.55),
                    "Never build notes from a raw book alone. Build from demands.")


def build_slide_04(prs):
    """Layer 2 — Micro-Thematic Structuring"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "L A Y E R  2", font_size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "Micro-Thematic Structuring", font_size=34, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.65), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.85), Inches(6.8), Inches(0.5),
                "Don't build notes topic-by-topic. Build them demand-by-demand.",
                font_size=15, color=GOLD, font_name=TITLE_FONT)
    # Before card
    add_card(slide, Inches(0.8), Inches(2.6), Inches(3.0), Inches(2.8), border_color=CARD_BORDER)
    add_multiline(slide, Inches(1.0), Inches(2.7), Inches(2.7), Inches(2.6), [
        ("BEFORE", MUTED, True, 11),
        ("", BODY_TEXT, False, 6),
        ("22 PYQs on Globalisation", BODY_TEXT, False, 12),
        ("Impact on women", BODY_TEXT, False, 11),
        ("Impact on elderly", BODY_TEXT, False, 11),
        ("Cultural identity", BODY_TEXT, False, 11),
        ("Technology effects", BODY_TEXT, False, 11),
        ("Trade dynamics", BODY_TEXT, False, 11),
        ("...and 17 more overlaps", MUTED, False, 11),
    ])
    # Arrow
    add_textbox(slide, Inches(3.9), Inches(3.8), Inches(0.6), Inches(0.5),
                "→", font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # After card
    add_card(slide, Inches(4.5), Inches(2.6), Inches(3.0), Inches(2.8), border_color=TEAL)
    add_multiline(slide, Inches(4.7), Inches(2.7), Inches(2.7), Inches(2.6), [
        ("AFTER", TEAL, True, 11),
        ("", BODY_TEXT, False, 6),
        ("9 Real Dimensions", GOLD, True, 15),
        ("", BODY_TEXT, False, 5),
        ("Redundancy removed", BODY_TEXT, False, 12),
        ("Same content coverage", BODY_TEXT, False, 12),
        ("1/3 the pages", BODY_TEXT, False, 12),
        ("3x retrieval speed", BODY_TEXT, False, 12),
    ])
    add_insider_tip(slide, Inches(0.8), Inches(5.8), Inches(6.8), Inches(0.6),
                    "UPSC rotates the same 3-4 dimensions in different costumes. Find the costume, ignore the noise.")



def build_slide_05(prs):
    """Layer 3 — Cross-Paper Consolidation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "L A Y E R  3", font_size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "Cross-Paper Consolidation", font_size=34, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.65), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.85), Inches(6.5), Inches(0.4),
                "One fact. Multiple papers. Maximum marks per unit of knowledge.", font_size=14, color=MUTED)
    # Theme label
    add_textbox(slide, Inches(2.5), Inches(2.4), Inches(3.0), Inches(0.4),
                "THEME: Women", font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # Four GS cards
    papers = [
        ("GS1", "Society: Status,\neducation, workforce"),
        ("GS2", "Governance: SHGs,\npolicies, rights"),
        ("GS3", "Economy: Labour\nparticipation, MUDRA"),
        ("GS4", "Ethics: Feminist\nethics, case studies"),
    ]
    for i, (gs, desc) in enumerate(papers):
        x = Inches(0.8 + i * 1.8)
        add_card(slide, x, Inches(3.0), Inches(1.6), Inches(2.0), border_color=CARD_BORDER)
        add_textbox(slide, x, Inches(3.1), Inches(1.6), Inches(0.35),
                    gs, font_size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_multiline(slide, int(x + Inches(0.1)), Inches(3.5), Inches(1.4), Inches(1.3),
                      [(line, BODY_TEXT, False, 10) for line in desc.split("\n")],
                      align=PP_ALIGN.CENTER)
    # Principle box
    add_card(slide, Inches(0.8), Inches(5.3), Inches(6.8), Inches(0.55), border_color=CARD_BORDER)
    add_multiline(slide, Inches(1.0), Inches(5.35), Inches(6.4), Inches(0.45), [
        ("Principle: Content is limited; application is unlimited.", GOLD, True, 13),
    ])
    add_insider_tip(slide, Inches(0.8), Inches(6.1), Inches(6.8), Inches(0.55),
                    "Ask for every fact: 'Where else could this live?' — one habit that shrinks total load.")


def build_slide_06(prs):
    """Layer 4 — Periodic Enrichment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "L A Y E R  4", font_size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "Periodic Enrichment", font_size=38, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.7), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(6.8), Inches(0.4),
                "Fold new content into existing structure. Never start a separate pile.", font_size=14, color=MUTED)
    # Three stage cards with arrows
    stages = [
        ("Month 1–3", "Build base structure\nfrom syllabus + PYQs"),
        ("Month 4–8", "Fold current affairs\ninto existing themes"),
        ("Month 9–12", "Final enrichment +\nrevision-ready polish"),
    ]
    for i, (period, desc) in enumerate(stages):
        x = Inches(0.8 + i * 2.4)
        add_card(slide, x, Inches(2.5), Inches(2.1), Inches(2.2), border_color=CARD_BORDER)
        add_textbox(slide, x, Inches(2.65), Inches(2.1), Inches(0.4),
                    period, font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_multiline(slide, int(x + Inches(0.15)), Inches(3.2), Inches(1.8), Inches(1.2),
                      [(line, BODY_TEXT, False, 11) for line in desc.split("\n")],
                      align=PP_ALIGN.CENTER)
    # Arrows
    add_textbox(slide, Inches(3.0), Inches(3.4), Inches(0.5), Inches(0.4),
                "→", font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.4), Inches(3.4), Inches(0.5), Inches(0.4),
                "→", font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # Rule box
    add_card(slide, Inches(0.8), Inches(5.1), Inches(6.8), Inches(0.55), border_color=GOLD)
    add_multiline(slide, Inches(1.0), Inches(5.15), Inches(6.4), Inches(0.45), [
        ("Rule: New information has an obvious home the moment you find it.", GOLD, True, 13),
    ])
    add_insider_tip(slide, Inches(0.8), Inches(5.9), Inches(6.8), Inches(0.55),
                    "Structure-first, content-second. The app doesn't save you if the folder logic is chaotic.")



def build_slide_07(prs):
    """Layer 5 — Reflection in Answer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "L A Y E R  5", font_size=12, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "Reflection in Answer", font_size=38, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.7), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(6.8), Inches(0.5),
                "Your copy is a mirror of your notes — nothing more, nothing less.",
                font_size=15, color=GOLD, font_name=TITLE_FONT)
    # Left: Note column
    add_card(slide, Inches(0.8), Inches(2.6), Inches(3.0), Inches(2.8), border_color=CARD_BORDER)
    add_multiline(slide, Inches(1.0), Inches(2.7), Inches(2.7), Inches(2.6), [
        ("YOUR NOTE", MUTED, True, 11),
        ("", BODY_TEXT, False, 6),
        ("GHI: 107 (2023)", BODY_TEXT, False, 12),
        ("India rank: 111/125", BODY_TEXT, False, 12),
        ("Child wasting: 18.7%", BODY_TEXT, False, 12),
        ("ICDS + PM Poshan link", BODY_TEXT, False, 12),
        ("", BODY_TEXT, False, 6),
        ("Keyword: nutritional", GOLD, False, 11),
        ("sovereignty", GOLD, False, 11),
    ])
    # Arrow
    add_textbox(slide, Inches(3.9), Inches(3.8), Inches(0.6), Inches(0.5),
                "→", font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # Right: Answer column
    add_card(slide, Inches(4.5), Inches(2.6), Inches(3.0), Inches(2.8), border_color=TEAL)
    add_multiline(slide, Inches(4.7), Inches(2.7), Inches(2.7), Inches(2.6), [
        ("YOUR ANSWER", TEAL, True, 11),
        ("", BODY_TEXT, False, 6),
        ("India's GHI score of 107", BODY_TEXT, False, 12),
        ("(rank 111/125) reveals a", BODY_TEXT, False, 12),
        ("systemic crisis in nutritional", BODY_TEXT, False, 12),
        ("sovereignty. With child", BODY_TEXT, False, 12),
        ("wasting at 18.7%, convergence", BODY_TEXT, False, 12),
        ("of ICDS and PM Poshan", BODY_TEXT, False, 12),
        ("remains critical...", BODY_TEXT, False, 12),
    ])
    add_insider_tip(slide, Inches(0.8), Inches(5.8), Inches(6.8), Inches(0.6),
                    "Model answers are already in exam-language. Build from them to skip the translation lag.")


def build_slide_08(prs):
    """The One-Page Rule"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "S T A T I C  R U L E", font_size=12, bold=True, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "The One-Page Rule", font_size=40, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.7), Inches(2.5))
    # Big "1"
    add_textbox(slide, Inches(2.5), Inches(2.0), Inches(3.0), Inches(2.5),
                "1", font_size=120, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font_name=TITLE_FONT)
    # Rule box
    add_card(slide, Inches(0.8), Inches(4.4), Inches(6.8), Inches(1.4), border_color=CARD_BORDER)
    add_multiline(slide, Inches(1.0), Inches(4.5), Inches(6.4), Inches(1.2), [
        ("If it doesn't fit on one page, it isn't a note yet.", GOLD, True, 15),
        ("It's still raw material.", MUTED, False, 13),
        ("", BODY_TEXT, False, 5),
        ("One page forces demand-deconstruction. You can't cheat", BODY_TEXT, False, 12),
        ("your way to one page by writing smaller — only by removing overlap.", BODY_TEXT, False, 12),
    ])
    add_insider_tip(slide, Inches(0.8), Inches(6.1), Inches(6.8), Inches(0.55),
                    "The one-page constraint is a forcing function, not an arbitrary limit.")



def build_slide_09(prs):
    """Digital vs Physical"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "T O O L  S T R A T E G Y", font_size=12, bold=True, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "Digital vs Physical", font_size=38, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.7), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(6.5), Inches(0.4),
                "Different content types need different tools.", font_size=14, color=MUTED)
    # Left — Handwritten
    add_card(slide, Inches(0.8), Inches(2.5), Inches(3.2), Inches(3.2), border_color=GOLD)
    add_multiline(slide, Inches(1.0), Inches(2.6), Inches(2.9), Inches(3.0), [
        ("HANDWRITTEN", GOLD, True, 13),
        ("Static Subjects", MUTED, False, 11),
        ("", BODY_TEXT, False, 6),
        ("1. History maps + timelines", BODY_TEXT, False, 12),
        ("2. Polity fundamentals", BODY_TEXT, False, 12),
        ("3. Geography one-pagers", BODY_TEXT, False, 12),
        ("4. Ethics case templates", BODY_TEXT, False, 12),
        ("", BODY_TEXT, False, 6),
        ("Why: Forces pre-processing", TEAL, False, 11),
        ("that typing skips entirely.", TEAL, False, 11),
    ])
    # Right — Digital
    add_card(slide, Inches(4.3), Inches(2.5), Inches(3.2), Inches(3.2), border_color=TEAL)
    add_multiline(slide, Inches(4.5), Inches(2.6), Inches(2.9), Inches(3.0), [
        ("DIGITAL", TEAL, True, 13),
        ("Current Affairs", MUTED, False, 11),
        ("", BODY_TEXT, False, 6),
        ("1. Screenshot key data", BODY_TEXT, False, 12),
        ("2. Tag by GS paper + theme", BODY_TEXT, False, 12),
        ("3. Update as events evolve", BODY_TEXT, False, 12),
        ("4. Search instantly", BODY_TEXT, False, 12),
        ("", BODY_TEXT, False, 6),
        ("Why: Fastest-decaying content", GOLD, False, 11),
        ("needs fastest tool.", GOLD, False, 11),
    ])
    add_insider_tip(slide, Inches(0.8), Inches(6.1), Inches(6.8), Inches(0.6),
                    "Hand-copying current affairs = using your slowest tool on your fastest-decaying content.")


def build_slide_10(prs):
    """AI's Two Jobs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "T H E  A I  E D G E", font_size=12, bold=True, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "AI's Two Jobs", font_size=38, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.7), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(6.8), Inches(0.4),
                "Velocity on pattern-recognition. Not a replacement for judgment.", font_size=14, color=MUTED)
    # Job 1
    add_card(slide, Inches(0.8), Inches(2.5), Inches(3.2), Inches(2.8), border_color=CARD_BORDER)
    add_textbox(slide, Inches(0.8), Inches(2.6), Inches(3.2), Inches(0.5),
                "1", font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font_name=TITLE_FONT)
    add_multiline(slide, Inches(1.0), Inches(3.2), Inches(2.9), Inches(1.8), [
        ("Cluster PYQs Into Themes", GOLD, True, 13),
        ("", BODY_TEXT, False, 5),
        ("Feed it a decade of PYQs.", BODY_TEXT, False, 12),
        ("It finds recurring patterns", BODY_TEXT, False, 12),
        ("in minutes — replacing a", BODY_TEXT, False, 12),
        ("weekend of manual work.", BODY_TEXT, False, 12),
    ], align=PP_ALIGN.CENTER)
    # Job 2
    add_card(slide, Inches(4.3), Inches(2.5), Inches(3.2), Inches(2.8), border_color=CARD_BORDER)
    add_textbox(slide, Inches(4.3), Inches(2.6), Inches(3.2), Inches(0.5),
                "2", font_size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font_name=TITLE_FONT)
    add_multiline(slide, Inches(4.5), Inches(3.2), Inches(2.9), Inches(1.8), [
        ("Draft Practice Case Studies", GOLD, True, 13),
        ("", BODY_TEXT, False, 5),
        ("Generate ethics dilemmas", BODY_TEXT, False, 12),
        ("around a keyword. Select", BODY_TEXT, False, 12),
        ("and refine the 5-6 that", BODY_TEXT, False, 12),
        ("fit your answer templates.", BODY_TEXT, False, 12),
    ], align=PP_ALIGN.CENTER)
    # Bottom note
    add_card(slide, Inches(0.8), Inches(5.6), Inches(6.8), Inches(0.45), border_color=CARD_BORDER)
    add_textbox(slide, Inches(1.0), Inches(5.63), Inches(6.4), Inches(0.4),
                "AI removes the tedious first pass. Your judgment decides what stays.",
                font_size=12, color=BODY_TEXT, align=PP_ALIGN.CENTER)
    add_insider_tip(slide, Inches(0.8), Inches(6.25), Inches(6.8), Inches(0.5),
                    "Human-in-the-loop always. AI accelerates — it does not replace thinking.")



def build_slide_11(prs):
    """Mains vs Prelims — Tactical Split"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "T A C T I C A L  S P L I T", font_size=12, bold=True, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "Mains vs Prelims", font_size=38, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.7), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(6.5), Inches(0.4),
                "Same notes. Two different deployment modes.", font_size=14, color=MUTED)
    # MAINS card
    add_card(slide, Inches(0.8), Inches(2.5), Inches(3.2), Inches(3.3), border_color=GOLD)
    add_multiline(slide, Inches(1.0), Inches(2.6), Inches(2.9), Inches(3.1), [
        ("MAINS", GOLD, True, 14),
        ("Argument Synthesis", MUTED, False, 11),
        ("", BODY_TEXT, False, 6),
        ("1. 250-word ceiling constraint", BODY_TEXT, False, 12),
        ("2. Sub-150 words per dimension", BODY_TEXT, False, 12),
        ("3. Intro-Body-Conclusion pre-built", BODY_TEXT, False, 12),
        ("4. Skeleton reusable across topics", BODY_TEXT, False, 12),
        ("", BODY_TEXT, False, 6),
        ("Goal: Reproduce structured", TEAL, False, 11),
        ("argument from one-page note", TEAL, False, 11),
        ("in under 7 minutes.", TEAL, False, 11),
    ])
    # PRELIMS card
    add_card(slide, Inches(4.3), Inches(2.5), Inches(3.2), Inches(3.3), border_color=TEAL)
    add_multiline(slide, Inches(4.5), Inches(2.6), Inches(2.9), Inches(3.1), [
        ("PRELIMS", TEAL, True, 14),
        ("Fact Recall + Clarity", MUTED, False, 11),
        ("", BODY_TEXT, False, 6),
        ("1. Pure fact-retention", BODY_TEXT, False, 12),
        ("2. Conceptual clarity", BODY_TEXT, False, 12),
        ("3. Recognition speed", BODY_TEXT, False, 12),
        ("4. Same map, dual payoff", BODY_TEXT, False, 12),
        ("", BODY_TEXT, False, 6),
        ("Goal: Recognise correct option", GOLD, False, 11),
        ("in under 60 seconds using", GOLD, False, 11),
        ("the same integrated map.", GOLD, False, 11),
    ])
    add_insider_tip(slide, Inches(0.8), Inches(6.1), Inches(6.8), Inches(0.55),
                    "Integrated notes = one build, two payoffs. Don't maintain separate systems.")


def build_slide_12(prs):
    """UKPSC-UPSC Overlap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "U K P S C  S T R A T E G Y", font_size=12, bold=True, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "The UKPSC–UPSC Overlap", font_size=36, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.7), Inches(2.5))
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(6.8), Inches(0.4),
                "Your efficiency multiplier if preparing for both exams.", font_size=14, color=MUTED)
    # Center: Shared core
    add_card(slide, Inches(1.8), Inches(2.5), Inches(4.4), Inches(1.8), border_color=GOLD)
    add_multiline(slide, Inches(2.0), Inches(2.6), Inches(4.0), Inches(1.6), [
        ("SHARED STATIC CORE", GOLD, True, 14),
        ("", BODY_TEXT, False, 5),
        ("History  |  Geography  |  Polity", BODY_TEXT, False, 13),
        ("", BODY_TEXT, False, 4),
        ("Build ONE integrated map.", BODY_TEXT, False, 12),
        ("Barely diverges between exams.", BODY_TEXT, False, 12),
    ], align=PP_ALIGN.CENTER)
    # Left: UKPSC layer
    add_card(slide, Inches(0.8), Inches(4.7), Inches(3.2), Inches(1.5), border_color=TEAL)
    add_multiline(slide, Inches(1.0), Inches(4.8), Inches(2.9), Inches(1.3), [
        ("+ UKPSC LAYER", TEAL, True, 12),
        ("", BODY_TEXT, False, 4),
        ("State-specific current affairs", BODY_TEXT, False, 11),
        ("Uttarakhand schemes + movements", BODY_TEXT, False, 11),
        ("State geography specifics", BODY_TEXT, False, 11),
    ], align=PP_ALIGN.CENTER)
    # Right: UPSC layer
    add_card(slide, Inches(4.3), Inches(4.7), Inches(3.2), Inches(1.5), border_color=GOLD)
    add_multiline(slide, Inches(4.5), Inches(4.8), Inches(2.9), Inches(1.3), [
        ("+ UPSC LAYER", GOLD, True, 12),
        ("", BODY_TEXT, False, 4),
        ("National current affairs", BODY_TEXT, False, 11),
        ("International relations", BODY_TEXT, False, 11),
        ("Advanced governance topics", BODY_TEXT, False, 11),
    ], align=PP_ALIGN.CENTER)
    add_insider_tip(slide, Inches(0.8), Inches(6.5), Inches(6.8), Inches(0.5),
                    "Don't build two separate systems. Layer state-specific on top of a shared foundation.")



def build_slide_13(prs):
    """The 4 Topper Secrets — Recap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6.5), Inches(0.4),
                "T O P P E R  S E C R E T S", font_size=12, bold=True, color=MUTED)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(6.5), Inches(0.7),
                "The 4 Non-Obvious Truths", font_size=36, bold=True, color=BODY_TEXT, font_name=TITLE_FONT)
    add_gold_line(slide, Inches(0.8), Inches(1.65), Inches(2.5))
    # 4 numbered secrets
    secrets = [
        ("1", "Demand Deconstruction", "Build notes by demand, not by chapter. 22 PYQ variations collapse into 8-10 reusable dimensions."),
        ("2", "Model Answers Over Books", "Build from exam-language sources. Skip the book-to-answer translation lag entirely."),
        ("3", "Skip What You Know", "Notes = record of what you'd forget + phrases that differentiate you from everyone else."),
        ("4", "Cross-Paper Transfer", "Every fact should earn marks in multiple papers. Ask: where else could this live?"),
    ]
    for i, (num, title, desc) in enumerate(secrets):
        y = Inches(2.0 + i * 1.05)
        # Gold number box
        add_card(slide, Inches(0.8), y, Inches(0.5), Inches(0.5), fill=GOLD)
        add_textbox(slide, Inches(0.8), y, Inches(0.5), Inches(0.5),
                    num, font_size=16, bold=True, color=CARBON, align=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, Inches(1.5), y, Inches(6.0), Inches(0.4),
                    title, font_size=15, bold=True, color=BODY_TEXT)
        # Description
        add_textbox(slide, Inches(1.5), int(y + Inches(0.4)), Inches(6.0), Inches(0.5),
                    desc, font_size=11, color=MUTED)
    add_insider_tip(slide, Inches(0.8), Inches(6.3), Inches(6.8), Inches(0.55),
                    "These secrets compound. Each layer makes the next one exponentially more effective.")


def build_slide_14(prs):
    """Closing / Sign-Off"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # Centered minimal
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(6.8), Inches(0.4),
                "T H E  S I N G L E  T A K E A W A Y", font_size=12, bold=True, color=MUTED,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(6.8), Inches(0.8),
                "Build the demand map", font_size=36, bold=True, color=BODY_TEXT,
                align=PP_ALIGN.CENTER, font_name=TITLE_FONT)
    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(6.8), Inches(0.8),
                "before you build the note.", font_size=36, bold=True, color=BODY_TEXT,
                align=PP_ALIGN.CENTER, font_name=TITLE_FONT)
    # Gold separator
    add_gold_line(slide, Inches(3.0), Inches(4.2), Inches(2.0))
    # Sub-text
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(6.8), Inches(0.5),
                "That's the whole secret.", font_size=16, color=MUTED, align=PP_ALIGN.CENTER)
    # Brand
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(6.8), Inches(0.4),
                "UKPSC DECODED", font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(6.8), Inches(0.4),
                "Prepare Smarter. Not Longer.", font_size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# MAIN
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Build all 14 slides
    build_slide_01(prs)
    build_slide_02(prs)
    build_slide_03(prs)
    build_slide_04(prs)
    build_slide_05(prs)
    build_slide_06(prs)
    build_slide_07(prs)
    build_slide_08(prs)
    build_slide_09(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    build_slide_12(prs)
    build_slide_13(prs)
    build_slide_14(prs)

    out_path = "/projects/sandbox/Ukpscdecoded/UKPSC_Video2_NoteMaking_Slides.pptx"
    prs.save(out_path)
    print(f"Generated: {out_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("Ready for Keynote import!")


if __name__ == "__main__":
    main()
